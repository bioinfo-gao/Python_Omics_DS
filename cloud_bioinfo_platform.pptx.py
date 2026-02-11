from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Function to set background color of slide
def set_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Slide 1: Title slide
slide_layout = prs.slide_layouts[0]  # 0 is the layout for title slide
def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle
    
    # Set background color to light blue
    set_background_color(slide, RGBColor(173, 216, 230))

add_title_slide(prs, "Cloud Based Bioinformation Platform构想", "A Vision for Future Bioinformatics")

# Slide 2: Table of Contents
slide_layout = prs.slide_layouts[1]  # 1 is the layout for section header
def add_table_of_contents(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Table of Contents"
    content_shape.text = (
        "1. Introduction\n"
        "2. Platform Architecture Design\n"
        "3. Functional Modules Detailed Planning\n"
        "4. Conclusion"
    )
    
    # Set background color to light blue
    set_background_color(slide, RGBColor(173, 216, 230))

add_table_of_contents(prs)

# Slide 3: Introduction
slide_layout = prs.slide_layouts[1]
def add_introduction_slide(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Introduction"
    content_shape.text = (
        "The Cloud Based Bioinformation Platform aims to provide scalable, efficient, and collaborative solutions for bioinformatics research."
    )

add_introduction_slide(prs)

# Slide 4: Platform Architecture Design - Infrastructure Layer
slide_layout = prs.slide_layouts[1]
def add_infrastructure_layer_slide(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Platform Architecture Design - Infrastructure Layer"
    content_shape.text = (
        "Cloud Service Selection: AWS (Batch/Nextflow) or Google Cloud (Life Sciences API).\n"
        "Data Delivery: Raw data stored in AWS S3.\n"
        "CLI Tools and Web-based Resumable Downloads (for Gb-level data downloads).\n"
        "Permission Management: IAM Policies or Temporary Signed URLs to ensure data visibility only to clients."
    )

add_infrastructure_layer_slide(prs)

# Slide 5: Platform Architecture Design - Compute Engine Layer
slide_layout = prs.slide_layouts[1]
def add_compute_engine_layer_slide(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Platform Architecture Design - Compute Engine Layer"
    content_shape.text = (
        "Technology Stack: Nextflow or Snakemake for standardized pipelines.\n"
        "Bulk RNA-seq: fastp -> STAR -> featureCounts -> DESeq2.\n"
        "10x Single Cell: CellRanger -> Seurat/Scanpy.\n"
        "Containerization: All software environments encapsulated in Docker for consistency across cloud instances."
    )

add_compute_engine_layer_slide(prs)

# Slide 6: Platform Architecture Design - AI Interaction Layer
slide_layout = prs.slide_layouts[1]
def add_ai_interaction_layer_slide(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Platform Architecture Design - AI Interaction Layer"
    content_shape.text = (
        "LLM Engine: Integrate Claude 3.5 Sonnet or GPT-4o for code generation and scientific chart understanding.\n"
        "Agent Mode: User input -> LLM generates Python (Scanpy/Seurat) or R code -> Runs in an isolated sandbox environment -> Returns result charts."
    )

add_ai_interaction_layer_slide(prs)

# Slide 7: Functional Modules Detailed Planning - Module One
slide_layout = prs.slide_layouts[1]
def add_module_one_slide(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Functional Modules Detailed Planning - Module One"
    content_shape.text = (
        "Raw Data Distribution (The 'Drop-box' for Bio-data):\n"
        "Client logs in to see a Dropbox-like interface to access Fastq files.\n"
        "Provide MD5 checksum to ensure transmission integrity."
    )

add_module_one_slide(prs)

# Slide 8: Functional Modules Detailed Planning - Module Two
slide_layout = prs.slide_layouts[1]
def add_module_two_slide(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Functional Modules Detailed Planning - Module Two"
    content_shape.text = (
        "Standard Analysis Pipeline (One-click Pipeline):\n"
        "User selects samples -> Select species -> Click 'Start Running'.\n"
        "Real-time progress bar (Log streaming)."
    )

add_module_two_slide(prs)

# Slide 9: Functional Modules Detailed Planning - Module Three
slide_layout = prs.slide_layouts[1]
def add_module_three_slide(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Functional Modules Detailed Planning - Module Three"
    content_shape.text = (
        "Interactive AI Laboratory (AI Bio-CoPilot):\n"
        "Predefined Toolset: Encapsulate commonly used R/Python plotting libraries (ggplot2, ComplexHeatmap).\n"
        "Natural Language Interaction: 'Change heatmap colors to red-green and cluster by gene family.'\n"
        "Instant Modification: All charts support PDF/SVG export for publication."
    )

add_module_three_slide(prs)

# Slide 10: Conclusion
slide_layout = prs.slide_layouts[1]
def add_conclusion_slide(prs):
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.shapes.placeholders[1]
    
    title_shape.text = "Conclusion"
    content_shape.text = (
        "The Cloud Based Bioinformation Platform represents a transformative approach to bioinformatics, offering unprecedented opportunities for innovation and discovery."
    )

add_conclusion_slide(prs)

# Save the presentation
prs.save('cloud_bioinfo_platform.pptx')